import re
import time
import os
import bcrypt
import random
import smtplib
from dotenv import load_dotenv
from flask import Flask, render_template, request, jsonify, redirect, session, flash, url_for
import google.generativeai as genai
from googleapiclient.discovery import build
from youtube_transcript_api import YouTubeTranscriptApi
from pymongo import MongoClient
from datetime import datetime
import secrets
import mysql.connector
import requests
import base64
import PyPDF2
from docx import Document
from pptx import Presentation
load_dotenv()

app = Flask(__name__)
app.secret_key = os.getenv("SECRET_KEY") or secrets.token_hex(32)

# Configure Gemini API
api_key=os.getenv("api_key")
genai.configure(api_key=api_key)

# Define Gemini model
model = genai.GenerativeModel("gemini-1.5-flash")

# Connect to MongoDB
mongo_url=os.getenv("mongo_url")
client = MongoClient(mongo_url)
db = client["chatbot_db"]
collection = db["chat_history"]

# Connect to MySQL (for user authentication)
def get_db_connection():
    connection = mysql.connector.connect(
        host=os.getenv("MYSQL_ADDON_HOST"),
        user=os.getenv("MYSQL_ADDON_USER"),
        password=os.getenv("MYSQL_ADDON_PASSWORD"),
        database=os.getenv("MYSQL_ADDON_DB"),
        port=int(os.getenv("MYSQL_ADDON_PORT", 3306))  # Default to 3306
    )
    return connection

youtube = build('youtube', 'v3', developerKey=os.getenv('YOUTUBE_API_KEY'))

def extract_video_id(url):
    """Extract video ID from various YouTube URL formats"""
    patterns = [
        r'(?:https?:\/\/)?(?:www\.)?youtu\.be\/([a-zA-Z0-9_-]+)',
        r'(?:https?:\/\/)?(?:www\.)?youtube\.com\/watch\?v=([a-zA-Z0-9_-]+)',
        r'(?:https?:\/\/)?(?:www\.)?youtube\.com\/embed\/([a-zA-Z0-9_-]+)',
        r'(?:https?:\/\/)?(?:www\.)?youtube\.com\/v\/([a-zA-Z0-9_-]+)',
        r'(?:https?:\/\/)?(?:www\.)?youtube\.com\/shorts\/([a-zA-Z0-9_-]+)'
    ]
    for pattern in patterns:
        match = re.search(pattern, url)
        if match:
            return match.group(1)
    return None

def get_video_details(video_id):
    """Get video details including title, channel, and statistics"""
    response = youtube.videos().list(
        part="snippet,statistics",
        id=video_id
    ).execute()
    
    if not response.get('items'):
        return None
        
    snippet = response['items'][0]['snippet']
    stats = response['items'][0]['statistics']
    
    return {
        'title': snippet['title'],
        'channel': snippet['channelTitle'],
        'views': stats.get('viewCount', 'N/A'),
        'likes': stats.get('likeCount', 'N/A'),
        'thumbnail': snippet['thumbnails']['high']['url'],
        'published': snippet['publishedAt'][:10]  # YYYY-MM-DD
    }
def get_transcript(video_id):
    """Improved transcript fetching with modern API handling"""
    try:
        from youtube_transcript_api import YouTubeTranscriptApi
        from youtube_transcript_api._errors import TranscriptsDisabled
        
        try:
            # Try to get English transcript first
            transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=['en'])
            return " ".join([entry['text'] for entry in transcript])
            
        except TranscriptsDisabled:
            # Try any available language if English fails
            transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
            transcript = transcript_list.find_transcript(
                [t.language_code for t in transcript_list]
            ).fetch()
            return " ".join([entry['text'] for entry in transcript])
            
    except Exception as e:
        print(f"Transcript error for {video_id}: {str(e)}")
        return None
    

def human_format(num):
    """Format large numbers"""
    num = int(num)
    if num >= 1_000_000:
        return f"{num/1_000_000:.1f}M"
    if num >= 1_000:
        return f"{num/1_000:.1f}K"
    return str(num)

@app.route("/process-youtube", methods=["POST"])
def process_youtube():
    if "username" not in session:
        return jsonify({"error": "User not logged in"}), 401

    try:
        data = request.get_json()
        if not data or "url" not in data:
            return jsonify({"error": "URL not provided"}), 400
            
        url = data["url"]
        video_id = extract_video_id(url)
        
        if not video_id:
            return jsonify({"error": "Invalid YouTube URL"}), 400

        # Initialize chat_key here for all paths
        chat_key = get_today_chat_key(session["username"])
        
        # Get video details
        video_details = get_video_details(video_id)
        if not video_details:
            return jsonify({"error": "Video not found"}), 404
                
        # Get transcript
        try:
            transcript = get_transcript(video_id)
            if not transcript:
                raise Exception("No transcript available")
                
            # Generate summary
            summary = generate_video_summary(transcript, video_details)
            
            chat_entry = {
                "user": url,
                "bot": summary,
                "code": [],
                "youtube_data": {
                    **video_details,
                    "video_id": video_id
                }
            }
            
            collection.update_one(
                {"_id": chat_key},
                {"$push": {"chat_history": chat_entry}},
                upsert=True
            )
            
            return jsonify({
                "success": True,
                "summary": summary,
                "youtube_data": video_details
            })
            
        except Exception as transcript_error:
            # Fallback analysis when transcript fails
            fallback_analysis = analyze_without_transcript(video_id, video_details)
            
            chat_entry = {
                "user": url,
                "bot": fallback_analysis,
                "code": [],
                "youtube_data": {
                    **video_details,
                    "video_id": video_id,
                    "no_transcript": True
                }
            }
            
            collection.update_one(
                {"_id": chat_key},
                {"$push": {"chat_history": chat_entry}},
                upsert=True
            )
            
            return jsonify({
                "success": False,
                "message": fallback_analysis,
                "youtube_data": video_details
            })
            
    except Exception as e:
        return jsonify({"error": f"Error processing video: {str(e)}"}), 500
    
def generate_video_summary(transcript, video_details):
    """Generate concise 5-10 line text summary"""
    prompt = f"""
    Create a short 5-10 line summary of this YouTube video in paragraph form. Include only:
    1. Content type (animation, tutorial, match highlights etc.)
    2. 2-3 key moments or topics
    3. Notable technical aspects (if relevant)
    4. View count and channel name
    5. Overall impression
    
    Keep it concise - no bullet points or section headers. Just 5-10 clear sentences.
    
    Video Details:
    Title: {video_details['title']}
    Channel: {video_details['channel']}
    Views: {video_details['views']}
    
    Transcript Excerpt:
    {transcript[:2000]}... [truncated if long]
    """
    
    try:
        response = model.generate_content(prompt)
        # Clean up the response
        summary = response.text.replace('\n', ' ').strip()
        return ' '.join(summary.split())  # Remove extra spaces
    except Exception as e:
        print(f"Summary error: {str(e)}")
        return default_fallback_summary(video_details)

def default_fallback_summary(video_details):
    """Fallback when analysis fails"""
    return (
        f"This {video_details['channel']} video titled '{video_details['title']}' "
        f"has {human_format(video_details['views'])} views. "
        "Basic details are available but full analysis couldn't be generated."
    )

def analyze_without_transcript(video_id, video_details):
    """Concise analysis when no transcript exists"""
    prompt = f"""
    In 5-7 sentences, summarize this YouTube video based only on its metadata:
    Title: {video_details['title']}
    Channel: {video_details['channel']}
    Views: {video_details['views']}
    Published: {video_details['published']}
    
    Provide only: content type, likely topics, and audience appeal.
    """
    response = model.generate_content(prompt)
    return response.text.strip()

def get_today_chat_key(username):
    """Generate today's date key for the user's chat document"""
    today_date = datetime.now().strftime("%Y-%m-%d")
    return f"chat on {today_date} - {username}"

# Inappropriate keywords filter
INAPPROPRIATE_KEYWORDS = ["adult", "porn", "sex", "violence", "drugs", "hate"]

def is_inappropriate(content):
    """Check if content contains inappropriate keywords."""
    content = content.lower()
    return any(keyword in content for keyword in INAPPROPRIATE_KEYWORDS)


def extract_code(response_text):
    """
    Extracts code blocks from the response and replaces them with placeholders.
    """
    code_blocks = re.findall(r"```(.*?)```", response_text, re.DOTALL)
    text_without_code = re.sub(r"```.*?```", "[CODE_BLOCK]", response_text, flags=re.DOTALL).strip()

    return text_without_code, code_blocks

# Weather API Route
@app.route("/get-weather", methods=["POST"])
def get_weather():
    data = request.json
    city = data.get("city")
    if not city:
        return jsonify({"error": "City not provided"}), 400

    api_key = os.getenv("OPENWEATHERMAP_API_KEY")  # Fetch API key from .env
    url = f"https://api.openweathermap.org/data/2.5/weather?q={city}&appid={api_key}&units=metric"

    try:
        response = requests.get(url)
        response.raise_for_status()
        weather_data = response.json()

        if weather_data.get("cod") != 200:
            return jsonify({"error": "City not found"}), 404

        # Extract relevant weather information
        temperature = weather_data["main"]["temp"]
        weather_description = weather_data["weather"][0]["description"]
        weather_message = f"The weather in {city} is {weather_description} with a temperature of {temperature}°C."

        return jsonify({"success": True, "message": weather_message})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# News API Route
@app.route("/get-news", methods=["GET"])
def get_news():
    api_key = os.getenv("NEWS_API_KEY")  # Fetch API key from .env
    url = f"https://newsapi.org/v2/top-headlines?country=us&apiKey={api_key}"

    try:
        response = requests.get(url)
        response.raise_for_status()
        news_data = response.json()

        if news_data.get("status") != "ok":
            return jsonify({"error": "Failed to fetch news"}), 500

        # Extract top 5 headlines
        articles = news_data["articles"][:5]
        news_message = "Here are the latest headlines:\n" + "\n".join([article["title"] for article in articles])

        return jsonify({"success": True, "message": news_message})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# Translation API Route
@app.route("/translate", methods=["POST"])
def translate_text():
    data = request.json
    text = data.get("text")
    target_language = data.get("target_language")

    if not text or not target_language:
        return jsonify({"error": "Text or target language not provided"}), 400

    url = f"https://api.mymemory.translated.net/get?q={text}&langpair=en|{target_language}"

    try:
        response = requests.get(url)
        response.raise_for_status()
        translation_data = response.json()

        if translation_data.get("responseStatus") != 200:
            return jsonify({"error": "Translation failed"}), 500

        translated_text = translation_data["responseData"]["translatedText"]
        return jsonify({"success": True, "translated_text": translated_text})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/upload-file", methods=["POST"])
def upload_file():
    if "file" not in request.files:
        return jsonify({"success": False, "error": "No file uploaded"}), 400

    file = request.files["file"]
    if file.filename == "":
        return jsonify({"success": False, "error": "No selected file"}), 400

    # Validate file type
    allowed_extensions = {"pdf", "docx", "pptx"}
    file_extension = file.filename.split(".")[-1].lower()
    if file_extension not in allowed_extensions:
        return jsonify({"success": False, "error": "Invalid file type"}), 400

    # Save the file temporarily
    file_path = os.path.join("uploads", file.filename)
    file.save(file_path)

    # Extract text based on file type
    if file_extension == "pdf":
        text = extract_text_from_pdf(file_path)
    elif file_extension == "docx":
        text = extract_text_from_docx(file_path)
    elif file_extension == "pptx":
        text = extract_text_from_pptx(file_path)
    else:
        return jsonify({"success": False, "error": "Unsupported file type"}), 400

    # Clean up: Delete the temporary file
    os.remove(file_path)

    if not text:
        return jsonify({"success": False, "error": "Failed to extract text"}), 500

    # Summarize the text using Gemini API
    summary = summarize_text_with_gemini(text)
    if not summary:
        return jsonify({"success": False, "error": "Failed to summarize text"}), 500

    # Store the file name and summary in MongoDB
    username = session.get("username")
    if not username:
        return jsonify({"success": False, "error": "User not logged in"}), 401

    chat_key = get_today_chat_key(username)
    user_chat = collection.find_one({"_id": chat_key})
    if not user_chat:
        user_chat = {"_id": chat_key, "username": username, "chat_history": []}
        collection.insert_one(user_chat)

    # Add the file request and summary to the chat history
    chat_entry = {
        "user": file.filename,  # Store the file name as the user message
        "bot": summary,         # Store the summary as the bot message
        "code": [],             # No code blocks for file requests
    }
    user_chat["chat_history"].append(chat_entry)
    collection.update_one({"_id": chat_key}, {"$set": {"chat_history": user_chat["chat_history"]}})

    return jsonify({
        "success": True,
        "message": "File processed successfully",
        "summary": summary,
    })
def extract_text_from_pdf(file_path):
    """Extract text from a PDF file."""
    try:
        with open(file_path, "rb") as file:
            reader = PyPDF2.PdfReader(file)
            text = ""
            for page in reader.pages:
                text += page.extract_text()
            return text
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        return None

def extract_text_from_docx(file_path):
    """Extract text from a DOCX file."""
    try:
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
        return None

def extract_text_from_pptx(file_path):
    """Extract text from a PPTX file."""
    try:
        ppt = Presentation(file_path)
        text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
        return None
    
def summarize_text_with_gemini(text):
    """Summarize text using the Gemini API."""
    try:
        prompt = f"Summarize the following text and provide the key points in them:\n{text}"
        response = model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        print(f"Error summarizing text with Gemini: {e}")
        return None
    
# Clarifai API credentials
CLARIFAI_API_KEY =os.getenv("CLARIFAI_API_KEY")
CLARIFAI_MODEL_URL = os.getenv("CLARIFAI_MODEL_URL")

def analyze_image_with_clarifai(image_path):
    """Send the image to Clarifai API and return the analysis results."""
    with open(image_path, "rb") as image_file:
        # Encode the image as base64
        image_data = base64.b64encode(image_file.read()).decode("utf-8")

    headers = {
        "Authorization": f"Key {CLARIFAI_API_KEY}",
        "Content-Type": "application/json",
    }

    payload = {
        "inputs": [
            {
                "data": {
                    "image": {
                        "base64": image_data,
                    }
                }
            }
        ]
    }

    try:
        response = requests.post(CLARIFAI_MODEL_URL, headers=headers, json=payload)
        response.raise_for_status()
        results = response.json()

        # Extract concepts (labels and confidence scores) from the response
        concepts = results["outputs"][0]["data"]["concepts"]
        labels_with_scores = [(concept["name"], concept["value"]) for concept in concepts]
        return labels_with_scores
    except Exception as e:
        print(f"Error analyzing image with Clarifai: {e}")
        return None

def generate_paragraph_with_gemini(labels_with_scores):
    """Send the Clarifai response to Gemini API and generate a descriptive paragraph."""
    # Format the labels and scores as a string
    labels_text = ", ".join([f"{label} ({score * 100:.2f}%)" for label, score in labels_with_scores])

    # Create a prompt for Gemini
    prompt = f"Describe the following image contents in a paragraph: {labels_text}"

    try:
        # Generate a response using Gemini API
        response = model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        print(f"Error generating paragraph with Gemini: {e}")
        return None

@app.route("/upload-image", methods=["POST"])
def upload_image():
    if "image" not in request.files:
        return jsonify({"success": False, "error": "No image uploaded"}), 400

    image = request.files["image"]
    if image.filename == "":
        return jsonify({"success": False, "error": "No selected file"}), 400

    # Validate file type
    allowed_extensions = {"png", "jpg", "jpeg"}
    if "." in image.filename and image.filename.split(".")[-1].lower() in allowed_extensions:
        # Save the image temporarily
        image_path = os.path.join("uploads", image.filename)
        image.save(image_path)

        # Analyze the image using Clarifai API
        labels_with_scores = analyze_image_with_clarifai(image_path)

        # Clean up: Delete the temporary image file
        os.remove(image_path)

        if labels_with_scores:
            # Send the Clarifai response to Gemini API
            paragraph = generate_paragraph_with_gemini(labels_with_scores)

            if paragraph:
                # Store the file name and Gemini response in MongoDB
                username = session.get("username")
                if not username:
                    return jsonify({"success": False, "error": "User not logged in"}), 401

                chat_key = get_today_chat_key(username)
                user_chat = collection.find_one({"_id": chat_key})
                if not user_chat:
                    user_chat = {"_id": chat_key, "username": username, "chat_history": []}
                    collection.insert_one(user_chat)

                # Add the image request and Gemini response to the chat history
                chat_entry = {
                    "user": image.filename,  # Store the file name as the user message
                    "bot": paragraph,       # Store the Gemini response as the bot message
                    "code": [],              # No code blocks for image requests
                }
                user_chat["chat_history"].append(chat_entry)
                collection.update_one({"_id": chat_key}, {"$set": {"chat_history": user_chat["chat_history"]}})

                return jsonify({
                    "success": True,
                    "message": "Image analyzed successfully",
                    "paragraph": paragraph,
                })
            else:
                return jsonify({"success": False, "error": "Failed to generate paragraph"}), 500
        else:
            return jsonify({"success": False, "error": "Failed to analyze image"}), 500
    else:
        return jsonify({"success": False, "error": "Invalid file type"}), 400
    
        
def get_today_chat_key(username):
    """Generate today's date key for the user's chat document."""
    today_date = datetime.now().strftime("%Y-%m-%d")
    return f"chat on {today_date} - {username}"

def send_otp(email, otp):
    sender_email = os.getenv("MAIL_USERNAME")
    sender_password = os.getenv("MAIL_PASSWORD")
    subject = "Welcome to Mohan's Mini Chatbot"
    body = f"Welcome to Mohan's Mini Chatbot\nYour OTP for password reset is: {otp}\nThanks for registering!"
    message = f"Subject: {subject}\n\n{body}"
    try:
        server = smtplib.SMTP("smtp.gmail.com", 587)
        server.starttls()
        server.login(sender_email, sender_password)
        server.sendmail(sender_email, email, message)
        server.quit()
        return True
    except Exception as e:
        print("Error sending email:", e)
        return False

def validate_password(password):
    return len(password) >= 8 and any(char.islower() for char in password) and any(char.isupper() for char in password)

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/signup', methods=['GET', 'POST'])
def signup():
    if request.method == 'POST':
        name = request.form['name']
        email = request.form['email']
        password = request.form['password']
        
        if not validate_password(password):
            flash("Password must contain at least one uppercase letter, one lowercase letter, and be at least 8 characters long.")
            return render_template('signup.html')

        connection = get_db_connection()
        cursor = connection.cursor()
        
        cursor.execute('SELECT * FROM users WHERE email = %s', (email,))
        existing_user = cursor.fetchone()
        
        if existing_user:
            flash("Email already exists! Please try a different one.")
            return render_template('signup.html')

        hashed_password = bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt())
        
        cursor.execute('INSERT INTO users (name, email, password) VALUES (%s, %s, %s)', (name, email, hashed_password.decode('utf-8')))
        connection.commit()
        
        cursor.close()
        connection.close()
        
        flash("Account successfully created! Please login.")
        return redirect(url_for('login'))  # Corrected: Use route function name
    return render_template('signup.html')
@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        email = request.form['email']
        password = request.form['password']
        
        connection = get_db_connection()
        cursor = connection.cursor()
        
        cursor.execute('SELECT name, email, password FROM users WHERE email = %s', (email,))
        user = cursor.fetchone()
        
        if user and bcrypt.checkpw(password.encode('utf-8'), user[2].encode('utf-8')):
            session["username"] = user[0]  # Store username in session
            return redirect(url_for('chatbot'))  # Corrected: Use route function name
        else:
            flash("Please enter correct details!", "error")
        
        cursor.close()
        connection.close()
    return render_template('login.html')

@app.route('/forgot-password', methods=['GET', 'POST'])
def forgot_password():
    if request.method == 'POST':
        email = request.form.get('email')
        otp = request.form.get('otp')
        new_password = request.form.get('new_password')
        confirm_password = request.form.get('confirm_password')

        connection = get_db_connection()
        cursor = connection.cursor()
        
        cursor.execute('SELECT * FROM users WHERE email = %s', (email,))
        user = cursor.fetchone()

        if email and not otp and not new_password:
            if user:
                session['reset_email'] = email
                session['otp'] = str(random.randint(100000, 999999))
                send_otp(email, session['otp'])
                flash("OTP sent successfully to your email!")
                return render_template('forgot_password.html', step=2)
            else:
                flash("No account exists with that email.")
                return render_template('forgot_password.html', step=1)

        elif otp:
            if otp == session.get('otp'):
                return render_template('forgot_password.html', step=3)
            else:
                flash("You have entered the wrong OTP. Please enter the correct OTP.")
                return render_template('forgot_password.html', step=2)

        elif new_password and confirm_password:
            if new_password != confirm_password:
                flash("Passwords do not match.")
                return render_template('forgot_password.html', step=3)

            if not validate_password(new_password):
                flash("Password must contain at least one uppercase letter, one lowercase letter, and be at least 8 characters long.")
                return render_template('forgot_password.html', step=3)

            hashed_password = bcrypt.hashpw(new_password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')
            cursor.execute('UPDATE users SET password = %s WHERE email = %s', (hashed_password, session['reset_email']))
            connection.commit()
            
            session.pop('reset_email', None)
            session.pop('otp', None)
            
            flash("Password updated successfully! Please login.")
            return redirect(url_for('login'))  # Corrected: Use route function name

        cursor.close()
        connection.close()
    return render_template('forgot_password.html', step=1)

@app.route("/get_history", methods=["GET"])
def get_history():
    if "username" not in session:
        return jsonify({"error": "User not logged in"}), 401

    username = session["username"]
    # Fetch all chat documents for the user
    user_chats = collection.find({"username": username})
    
    history = [{"chat_id": chat["_id"], "date": chat["_id"].split(" - ")[0]} for chat in user_chats]
    return jsonify({"history": history})


@app.route("/chatbot")
def chatbot():
    if "username" not in session:
        flash("You need to log in first.")
        return redirect(url_for('login'))  # Corrected: Use route function name
    return render_template("chatbot.html", username=session["username"])

@app.route("/logout")
def logout():
    """Logout user and clear session."""
    session.pop("username", None)
    return redirect("/")

@app.route("/load_chat_history/<chat_id>", methods=["GET"])
def load_chat_history(chat_id):
    if "username" not in session:
        return jsonify({"error": "User not logged in"}), 401

    # Fetch the chat document by chat_id
    user_chat = collection.find_one({"_id": chat_id})
    
    if user_chat:
        return jsonify({"chat_history": user_chat["chat_history"]})
    return jsonify({"error": "Chat history not found"})

@app.route("/get-username")
def get_username():
    username = session.get("username", "Guest")  # Default to "Guest" if no user is logged in
    return jsonify({"username": username})

@app.route("/chat", methods=["POST"])
def chat():
    if "username" not in session:
        return jsonify({"error": "User not logged in"}), 401

    username = session["username"]
    user_message = request.json.get("message")

    if not user_message:
        return jsonify({"error": "No message received"}), 400

    try:
        # Generate the key for today's chat document
        chat_key = get_today_chat_key(username)

        # Retrieve or create user chat history for today's date
        user_chat = collection.find_one({"_id": chat_key})
        if not user_chat:
            user_chat = {"_id": chat_key, "username": username, "chat_history": []}
            collection.insert_one(user_chat)

        chat_history = user_chat["chat_history"]
        chat_context = "\n".join([f"User: {chat['user']}\nAI: {chat['bot']}" for chat in chat_history])

        # Generate AI response
        context = chat_context + f"\nUser: {user_message}\nAI:"
        response = model.generate_content(context)

        # Process the response
        bot_reply = response.text.strip() if response else "I couldn't process that."
        bot_reply = bot_reply.replace("Bot:", "").replace("AI:", "").strip()

        # Extract code blocks and main text separately
        text_response, code_blocks = extract_code(bot_reply)

        # Store chat history in MongoDB
        chat_entry = {"user": user_message, "bot": text_response, "code": code_blocks}
        chat_history.append(chat_entry)
        collection.update_one({"_id": chat_key}, {"$set": {"chat_history": chat_history}})

        return jsonify({"reply": text_response, "code": code_blocks})

    except Exception as e:
        return jsonify({"reply": "An error occurred. Please try again."})

if __name__ == "__main__":
    app.run(debug=True)